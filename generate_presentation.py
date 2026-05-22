import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9 layout (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Modern Color Palette
    PRIMARY = RGBColor(30, 58, 138)    # Deep Royal Blue (#1e3a8a)
    SECONDARY = RGBColor(124, 58, 237) # Vibrant Purple (#7c3aed)
    DARK_TEXT = RGBColor(30, 41, 59)   # Slate Blue/Black (#1e293b)
    LIGHT_BG = RGBColor(248, 250, 252) # Clean Off-White (#f8fafc)
    MUTED_TEXT = RGBColor(100, 116, 139)# Muted Gray (#64748b)
    ACCENT = RGBColor(220, 38, 38)     # Deep Red (#dc2626)

    # Use a blank layout for complete custom design freedom
    blank_layout = prs.slide_layouts[6]
    
    # Helper to add standard header and footer to a slide
    def add_slide_decorations(slide, title_text):
        # Top banner/header
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = 'Inter'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        
        # Muted top category line
        p_sub = tf.add_paragraph()
        p_sub.text = "FILING BUDDY  |  PRACTICE MANAGEMENT SUITE"
        p_sub.font.name = 'Inter'
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = SECONDARY
        
        # Footer
        footerBox = slide.shapes.add_textbox(Inches(0.75), Inches(6.8), Inches(11.83), Inches(0.4))
        ftf = footerBox.text_frame
        fp = ftf.paragraphs[0]
        fp.text = "Confidential Client Presentation  |  Powered by Filing Buddy LLC"
        fp.font.name = 'Inter'
        fp.font.size = Pt(10)
        fp.font.color.rgb = MUTED_TEXT

    # ================= SLIDE 1: Title Slide =================
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Large centered text block for main title
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "FILING BUDDY"
    p1.font.name = 'Inter'
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "Modern Practice Management Suite for Accounting & Tax Firms"
    p2.font.name = 'Inter'
    p2.font.size = Pt(24)
    p2.font.color.rgb = SECONDARY
    
    p3 = tf.add_paragraph()
    p3.text = "\nStreamlining CA & Audit workflows, secure client portals, and dynamic role-based controls."
    p3.font.name = 'Inter'
    p3.font.size = Pt(16)
    p3.font.color.rgb = DARK_TEXT

    # ================= SLIDE 2: The Core Problem =================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide2, "The Challenge in Practice Management")
    
    txBox = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    points = [
        ("Complexity of Scale", "Managing hundreds of diverse compliance tasks across VAT, Corporate Tax, Audits, and Accounting leads to missed deadlines without central visibility."),
        ("Security Risks", "Sharing sensitive government portal login credentials across staff members via Excel sheets introduces immense security hazards and lack of accountability."),
        ("Operational Noise", "Staff members seeing all clients and operations causes information overload, leading to mistakes, delays, and critical communication gaps."),
        ("No Audit Trial", "Firms lack transparency in tracing when a return was filed, who handled it, or identifying backlog bottlenecks.")
    ]
    
    for title, desc in points:
        p_title = tf.add_paragraph()
        p_title.text = f"•  {title}"
        p_title.font.name = 'Inter'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = DARK_TEXT
        
        p_desc = tf.add_paragraph()
        p_desc.text = f"    {desc}\n"
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = MUTED_TEXT

    # ================= SLIDE 3: The Solution =================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide3, "The Solution: Filing Buddy")
    
    txBox = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    sol_points = [
        ("Unified Practice Dashboard", "One beautiful, responsive interface providing interactive task control, overdue notifications, and firm-wide compliance schedules."),
        ("Granular Role-Based Security", "Structured user workspaces ensuring staff members only access what they own, drastically lowering operational risk and data noise."),
        ("Government Portal Security Vault", "Secure, structured storage for trade licenses, portal login combinations, and primary contacts with access restricted on a need-to-know basis."),
        ("FTA Tracking & Notifications", "Dedicated tracking module for Federal Tax Authority submissions with automatic alert logs when additional queries are received.")
    ]
    
    for title, desc in sol_points:
        p_title = tf.add_paragraph()
        p_title.text = f"✓  {title}"
        p_title.font.name = 'Inter'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = PRIMARY
        
        p_desc = tf.add_paragraph()
        p_desc.text = f"    {desc}\n"
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = DARK_TEXT

    # ================= SLIDE 4: Today's High-Security Upgrades =================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide4, "Today's Core Performance Upgrades")
    
    txBox = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    upgrades = [
        ("1. Strict Client List Task Restriction (Task-Only Users)", 
         "• Staff members logged in under 'Task-Only' roles are restricted from seeing the full client directory.\n• They can ONLY see clients for whom work/tasks have been allotted to them, minimizing data noise and enhancing client privacy.\n• Programmatically implemented on the backend database level to ensure bulletproof data isolation."),
        ("2. Multi-Layer Delete Action Safeguards", 
         "• Deleting a client is a highly destructive operational risk.\n• Frontend: The delete button is completely hidden from managers and standard staff members, visible solely to Admins.\n• Backend: The deletion API endpoint is hard-blocked with JWT-role middleware requiring Admin role, returning 403 Forbidden for unauthorized requests.")
    ]
    
    for title, desc in upgrades:
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.name = 'Inter'
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = SECONDARY
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc + "\n"
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = DARK_TEXT

    # ================= SLIDE 5: Full Stack Architecture =================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide5, "Technical Architecture")
    
    txBox = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    arch = [
        ("Frontend Technology", "React.js, Vite for ultra-fast builds, custom Tailwind CSS styling, Lucide icons, and responsive tables with debounced searches."),
        ("Robust Backend Core", "Node.js with Express, RESTful APIs, JWT Authorization headers, robust role validation, and custom Mongoose model validation schemas."),
        ("Data Persistence Layer", "MongoDB (via Atlas for production, MongoMemoryServer in-memory DB for isolated, zero-dependency local testing)."),
        ("Live Cloud Infrastructure", "Vercel hosting the reactive frontend, Railway running the backend, with continuous automated CD integration via GitHub.")
    ]
    
    for title, desc in arch:
        p_title = tf.add_paragraph()
        p_title.text = f"❖  {title}"
        p_title.font.name = 'Inter'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = PRIMARY
        
        p_desc = tf.add_paragraph()
        p_desc.text = f"    {desc}\n"
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = DARK_TEXT

    # ================= SLIDE 6: Summary & Roadmap =================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide6, "Summary & Future Roadmap")
    
    txBox = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    roadmap = [
        ("Current System Status", "The system is fully verified end-to-end, extremely stable, and ready for high-level client presentations and operational deployment."),
        ("WhatsApp & Email Integrations", "Plan to implement automated SMS/WhatsApp alerts directly to clients for document updates, TRN requests, and payment remainders."),
        ("Document Generation Engine", "Dynamic PDF generation for client billing, VAT returns, and FTA correspondence letters directly from the dashboard.")
    ]
    
    for title, desc in roadmap:
        p_title = tf.add_paragraph()
        p_title.text = f"★  {title}"
        p_title.font.name = 'Inter'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = SECONDARY
        
        p_desc = tf.add_paragraph()
        p_desc.text = f"    {desc}\n"
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = DARK_TEXT

    # Save presentation
    prs.save("Filing_Buddy_Presentation.pptx")
    print("Presentation saved successfully as 'Filing_Buddy_Presentation.pptx'!")

if __name__ == '__main__':
    create_presentation()
